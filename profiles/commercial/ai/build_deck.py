"""Build AI Security lab PPT from talk-track content."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Microsoft-ish palette
NAVY = RGBColor(0x0B, 0x2D, 0x5B)
BLUE = RGBColor(0x00, 0x78, 0xD4)
TEAL = RGBColor(0x00, 0xB2, 0x94)
SLATE = RGBColor(0x32, 0x3B, 0x4C)
LIGHT = RGBColor(0xF3, 0xF2, 0xF1)
MUTED = RGBColor(0x60, 0x6F, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xFF, 0xA4, 0x4D)
RED = RGBColor(0xD1, 0x34, 0x38)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False
    return s


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=SLATE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(0.75)
    r.shadow.inherit = False
    return r


def add_rounded(slide, left, top, width, height, fill, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    r.adjustments[0] = 0.12
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(1)
    r.shadow.inherit = False
    return r


def add_header(slide, kicker, title, *, kicker_color=BLUE):
    # accent bar
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.15), Inches(0.5), kicker_color)
    add_text(slide, Inches(0.85), Inches(0.5), Inches(11), Inches(0.35),
             kicker.upper(), size=12, bold=True, color=kicker_color)
    add_text(slide, Inches(0.85), Inches(0.85), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    # hairline
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.02), LIGHT)


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "Integrated AI Security Lab  |  Purview + Copilot DLP + Shadow AI + Sentinel",
             size=9, color=MUTED)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.3),
             f"{idx} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, left, top, width, height, bullets, *, size=16, color=SLATE, bold=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "—  " + b
        r.font.name = "Segoe UI"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_table(slide, left, top, width, height, headers, rows, *, header_fill=NAVY,
              row_alt=LIGHT, header_text=WHITE, body_text=SLATE):
    cols = len(headers)
    total_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(total_rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    # headers
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h
        r.font.name = "Segoe UI"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = header_text
    # body
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else row_alt
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.name = "Segoe UI"
            r.font.size = Pt(11)
            r.font.color.rgb = body_text
    return tbl


# ---------- Slides ----------
slides_total_placeholder = [0]


def slide_title():
    s = add_slide()
    # left navy panel
    add_rect(s, 0, 0, Inches(5.0), prs.slide_height, NAVY)
    add_rect(s, Inches(0.6), Inches(2.6), Inches(3.8), Inches(0.08), TEAL)
    add_text(s, Inches(0.6), Inches(1.2), Inches(4.0), Inches(0.5),
             "PURVIEW LAB", size=13, bold=True, color=TEAL)
    add_text(s, Inches(0.6), Inches(1.6), Inches(4.2), Inches(1.2),
             "Integrated\nAI Security Lab", size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(3.0), Inches(4.2), Inches(2.0),
             "Copilot DLP +\nShadow AI prevention +\nSentinel SIEM,\none correlated loop.",
             size=16, color=WHITE)
    add_text(s, Inches(0.6), Inches(6.6), Inches(4.2), Inches(0.4),
             "PVAISec profile  |  commercial cloud", size=11, bold=True, color=TEAL)
    # right hero
    add_text(s, Inches(5.6), Inches(1.6), Inches(7.4), Inches(0.5),
             "THE EXEC LINE", size=13, bold=True, color=BLUE)
    add_text(s, Inches(5.6), Inches(2.0), Inches(7.4), Inches(2.5),
             "Four surfaces.\nOne policy model.\nOne SIEM pane.",
             size=36, bold=True, color=NAVY)
    add_text(s, Inches(5.6), Inches(4.6), Inches(7.4), Inches(1.8),
             "A complete Microsoft AI security posture:\nsanctioned AI guarded, unsanctioned AI prevented,\nInsider Risk adapting enforcement in real time,\nand every signal correlated in Sentinel.",
             size=16, color=SLATE)
    add_text(s, Inches(5.6), Inches(6.8), Inches(7.4), Inches(0.4),
             "Deploy: 20-25 min  |  Demo: 45-75 min hands-on",
             size=12, bold=True, color=MUTED)


def slide_problem():
    s = add_slide()
    add_header(s, "The problem", "Three AI security stories, usually treated as three projects")
    tiles = [
        ("Protect sanctioned AI",
         "Microsoft 365 Copilot.\nKeep it from leaking internal data or\nprocessing files it shouldn't.", BLUE),
        ("Prevent unsanctioned AI",
         "ChatGPT, Claude, Gemini.\nStop users pasting company data\ninto public AI sites.", TEAL),
        ("SIEM visibility",
         "Get signals to the SOC,\ncorrelate with the rest of security,\ndrive incident response.", SLATE),
    ]
    x = Inches(0.6)
    y = Inches(2.0)
    w = Inches(4.0)
    h = Inches(3.6)
    gap = Inches(0.15)
    for i, (title, body, color) in enumerate(tiles):
        left = Inches(0.6 + i * (4.0 + 0.15))
        add_rounded(s, left, y, w, h, WHITE, line=LIGHT)
        add_rect(s, left, y, w, Inches(0.5), color)
        add_text(s, left + Inches(0.3), y + Inches(0.08), w - Inches(0.6), Inches(0.5),
                 f"0{i+1}", size=14, bold=True, color=WHITE)
        add_text(s, left + Inches(0.3), y + Inches(0.7), w - Inches(0.6), Inches(0.6),
                 title, size=20, bold=True, color=NAVY)
        add_text(s, left + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(2.0),
                 body, size=14, color=SLATE)

    add_text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6),
             "The reality: these are one story. Signals from one surface feed the others.",
             size=18, bold=True, color=BLUE)


def slide_thesis():
    s = add_slide()
    add_header(s, "The thesis", "One system, not four tools", kicker_color=TEAL)
    add_text(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.0),
             "A user blocked pasting SSNs into ChatGPT is usually the same user\nflagged by Copilot prompt DLP — and their IRM score rises automatically,\ntightening DLP enforcement the next time.",
             size=20, color=SLATE)

    items = [
        ("Copilot DLP", "Sanctioned-AI guardrails at prompt + label level.", BLUE),
        ("Shadow AI", "Three-layer block at device, browser, network.", TEAL),
        ("Insider Risk", "Adaptive bridge — enforcement tiers respond to behavior.", AMBER),
        ("Sentinel", "Unified SIEM pane, cross-signal analytics rules.", NAVY),
    ]
    y = Inches(4.0)
    for i, (t, body, color) in enumerate(items):
        left = Inches(0.6 + i * 3.05)
        add_rounded(s, left, y, Inches(2.9), Inches(2.5), WHITE, line=LIGHT)
        add_rect(s, left, y, Inches(0.15), Inches(2.5), color)
        add_text(s, left + Inches(0.3), y + Inches(0.2), Inches(2.4), Inches(0.5),
                 t, size=16, bold=True, color=NAVY)
        add_text(s, left + Inches(0.3), y + Inches(0.85), Inches(2.4), Inches(1.5),
                 body, size=13, color=SLATE)


def slide_loop():
    s = add_slide()
    add_header(s, "The integrated loop", "Signals reinforce each other — automatically")
    steps = [
        ("1", "Copilot DLP block", "SSN/PHI in prompt → response blocked, no internal/web search"),
        ("2", "Shadow AI block", "Endpoint DLP intercepts paste to ChatGPT/Claude/Gemini"),
        ("3", "IRM score rises", "Risky AI Usage template escalates user tier"),
        ("4", "DLP tier tightens", "Same policy, stricter enforcement — no admin action"),
        ("5", "Sentinel correlates", "Cross-table rule fires on user with both signals in 4h"),
        ("6", "Auto-triage", "Logic App enriches incident with next-step comment"),
    ]
    y0 = Inches(1.9)
    box_w = Inches(3.95)
    box_h = Inches(1.55)
    for i, (n, t, body) in enumerate(steps):
        col = i % 3
        row = i // 3
        left = Inches(0.6 + col * 4.1)
        top = y0 + row * (box_h + Inches(0.25))
        add_rounded(s, left, top, box_w, box_h, WHITE, line=LIGHT)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.25),
                                  Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        circ.shadow.inherit = False
        add_text(s, left + Inches(0.2), top + Inches(0.3), Inches(0.7), Inches(0.6),
                 n, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(1.05), top + Inches(0.2), box_w - Inches(1.1), Inches(0.5),
                 t, size=15, bold=True, color=NAVY)
        add_text(s, left + Inches(1.05), top + Inches(0.65), box_w - Inches(1.1), Inches(1.0),
                 body, size=11, color=SLATE)

    add_text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
             "Entire loop closes inside Microsoft tooling. One config file deploys it all.",
             size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def slide_deploy_overview():
    s = add_slide()
    add_header(s, "What gets deployed", "One prefix: PVAISec")
    headers = ["Surface", "Count", "Resources"]
    rows = [
        ("Identity", "5 + 3", "Test users + security groups"),
        ("Sensitivity labels", "2 parents + 10 sublabels", "AI-specific taxonomy + 2 auto-label policies"),
        ("DLP policies", "5 (9 rules)", "Copilot Prompt, Copilot Label, Endpoint, Browser, Network"),
        ("Insider Risk", "3", "Risky AI usage, Data leaks, Departing users"),
        ("Retention", "5", "Exchange/SP/OD + Copilot + Enterprise AI + Other AI"),
        ("Comm Compliance", "2", "AI Activity Collection + PII/PHI Detection"),
        ("eDiscovery", "1", "Unified AI-Security-Incident-Review"),
        ("Conditional Access", "2 (report-only)", "Block high-risk AI + require MFA"),
        ("Sentinel", "7 rules + 2 workbooks + playbook", "Incl. 3 AI-specific + cross-signal correlation"),
        ("Test data", "4 emails + 5 docs", "Docs auto-labeled via Graph at deploy"),
    ]
    add_table(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.8), headers, rows)
    add_text(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.4),
             "Config: configs/commercial/ai-security-demo.json  |  Deploy: ~20-25 min  |  Propagation: 4h DLP, 60 min Sentinel",
             size=11, color=MUTED)


def slide_act1():
    s = add_slide()
    add_header(s, "Act 1  |  Discovery", "What AI activity is already happening in your tenant?")
    add_text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.5),
             "Three pre-configured audit searches", size=18, bold=True, color=NAVY)
    bullets = [
        "Copilot-Interaction-Audit — every Copilot use",
        "AI-DLP-Match-Audit — DLP rule matches across all surfaces",
        "AI-Policy-Override-Audit — attempts to override DLP blocks",
    ]
    add_bullets(s, Inches(0.6), Inches(2.6), Inches(6.0), Inches(2.5), bullets, size=15)

    add_text(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(0.5),
             "Then pivot to Sentinel workbook", size=18, bold=True, color=NAVY)
    add_text(s, Inches(7.0), Inches(2.6), Inches(6.0), Inches(3.0),
             "PVAISec-AI Risk Signals workbook:\n\n— Copilot DLP blocks over time\n— Shadow AI paste attempts by target site\n— Risky AI Usage IRM alerts\n— Cross-signal users (Copilot + IRM)",
             size=14, color=SLATE)

    add_rounded(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2), LIGHT)
    add_text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.0),
             "Cross-signal users are the key panel — the humans with both Copilot DLP blocks\nAND IRM AI scoring. Those are the real AI-risk people.",
             size=14, bold=True, color=NAVY)


def slide_act2_copilot():
    s = add_slide()
    add_header(s, "Act 2  |  Copilot DLP", "Sanctioned, but not unguarded")

    # Two policy cards
    y = Inches(1.95)
    h = Inches(3.9)
    w = Inches(5.9)

    # Card 1
    add_rounded(s, Inches(0.6), y, w, h, WHITE, line=LIGHT)
    add_rect(s, Inches(0.6), y, w, Inches(0.6), BLUE)
    add_text(s, Inches(0.8), y + Inches(0.12), w - Inches(0.3), Inches(0.5),
             "Copilot Prompt SIT Block   (preview)", size=16, bold=True, color=WHITE)
    add_text(s, Inches(0.8), y + Inches(0.85), w - Inches(0.3), Inches(0.4),
             "Location:  CopilotExperiences", size=12, bold=True, color=SLATE)
    add_text(s, Inches(0.8), y + Inches(1.3), w - Inches(0.3), Inches(0.4),
             "Condition:  SSN, Credit Card, PHI in prompt", size=12, color=SLATE)
    add_text(s, Inches(0.8), y + Inches(1.75), w - Inches(0.3), Inches(0.4),
             "Action:  Prevent Copilot from processing", size=12, color=SLATE)
    add_text(s, Inches(0.8), y + Inches(2.3), w - Inches(0.3), Inches(1.6),
             "One policy, three protections:\n— No Copilot response\n— No internal search\n— No web search\n\nSensitive string never leaves the guardrail.",
             size=12, color=SLATE)

    # Card 2
    left2 = Inches(0.6 + 5.9 + 0.3)
    add_rounded(s, left2, y, w, h, WHITE, line=LIGHT)
    add_rect(s, left2, y, w, Inches(0.6), TEAL)
    add_text(s, left2 + Inches(0.2), y + Inches(0.12), w - Inches(0.3), Inches(0.5),
             "Copilot Labeled Content Block   (GA)", size=16, bold=True, color=WHITE)
    add_text(s, left2 + Inches(0.2), y + Inches(0.85), w - Inches(0.3), Inches(0.4),
             "Location:  CopilotExperiences", size=12, bold=True, color=SLATE)
    add_text(s, left2 + Inches(0.2), y + Inches(1.3), w - Inches(0.3), Inches(0.4),
             "Condition:  label = AI Blocked or AI Regulated", size=12, color=SLATE)
    add_text(s, left2 + Inches(0.2), y + Inches(1.75), w - Inches(0.3), Inches(0.4),
             "Action:  Prevent Copilot from processing labeled file", size=12, color=SLATE)
    add_text(s, left2 + Inches(0.2), y + Inches(2.3), w - Inches(0.3), Inches(1.6),
             "Files auto-labeled on SSN detection\nimmediately become invisible to Copilot.\n\nUsers get a policy-driven message —\nnot a vague error.",
             size=12, color=SLATE)

    add_text(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.8),
             "Live demo beats:  \"Summarize benefits for 078-05-1120\" → blocked  |  \"Summarize Q4 Forecast\" → blocked (labeled)  |  \"What meetings this week?\" → normal",
             size=12, color=MUTED)


def slide_act3_shadow():
    s = add_slide()
    add_header(s, "Act 3  |  Shadow AI", "Three layers to close the unsanctioned AI gap")

    layers = [
        ("Layer 1  |  Devices", "Endpoint DLP",
         "Blocks paste/upload at the OS layer via\nDefender for Endpoint.\n\n10 AI sites in block list:\nChatGPT, Claude, Gemini, Perplexity, Poe,\nHuggingFace, DeepSeek, and more.", BLUE),
        ("Layer 2  |  Browser", "Edge for Business inline",
         "Inline inspection of prompt text in\nEdge for Business.\n\nCovers consumer AI sites:\nCopilot consumer, ChatGPT consumer,\nGemini, DeepSeek.", TEAL),
        ("Layer 3  |  Network", "SASE / SSE",
         "Network-layer DLP via SASE/SSE.\n\nCovers non-Edge browsers,\nnon-browser apps, APIs —\nanything that routes around the browser.", SLATE),
    ]
    y = Inches(1.95)
    h = Inches(4.2)
    w = Inches(3.95)
    for i, (title, sub, body, color) in enumerate(layers):
        left = Inches(0.6 + i * 4.1)
        add_rounded(s, left, y, w, h, WHITE, line=LIGHT)
        add_rect(s, left, y, w, Inches(0.7), color)
        add_text(s, left + Inches(0.3), y + Inches(0.1), w - Inches(0.3), Inches(0.5),
                 title, size=15, bold=True, color=WHITE)
        add_text(s, left + Inches(0.3), y + Inches(0.85), w - Inches(0.3), Inches(0.5),
                 sub, size=14, bold=True, color=NAVY)
        add_text(s, left + Inches(0.3), y + Inches(1.45), w - Inches(0.3), Inches(2.6),
                 body, size=12, color=SLATE)

    add_rounded(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.7), LIGHT)
    add_text(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.55),
             "The asymmetry:  Copilot DLP guides users to the sanctioned path.  External AI gets hard blocks.  That steers behavior.",
             size=13, bold=True, color=NAVY)


def slide_act4_irm():
    s = add_slide()
    add_header(s, "Act 4  |  Insider Risk", "The adaptive bridge")

    add_text(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(0.5),
             "Three IRM policies", size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(3.0), [
        "Risky AI Usage Watch — Copilot prompt injection, protected material access",
        "AI Data Exfiltration Watch — correlates DLP matches with data-leak signals",
        "Departing User AI Risk — elevates risk for departing users who touch AI",
    ], size=14)

    add_text(s, Inches(0.6), Inches(4.5), Inches(6.0), Inches(0.5),
             "Wizard defaults used", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.0), [
        "Scope: all users and groups (no priority scoping)",
        "Content: 1 random label + 1 SIT + 1 trainable classifier",
        "Detection: all indicators and triggering events selected",
    ], size=12)

    # Tier table on right
    add_text(s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(0.5),
             "Tiered enforcement  (Shadow AI Endpoint)", size=18, bold=True, color=NAVY)
    add_table(s, Inches(7.0), Inches(2.5), Inches(5.7), Inches(2.4),
              ["User risk", "DLP rule", "Enforcement"],
              [
                  ("Elevated", "Endpoint AI Block", "Hard block"),
                  ("Moderate", "Endpoint AI Warn", "Allow + justify"),
                  ("Minor", "Endpoint AI Audit", "Audit only"),
              ])
    add_rounded(s, Inches(7.0), Inches(5.1), Inches(5.7), Inches(1.9), LIGHT)
    add_text(s, Inches(7.2), Inches(5.25), Inches(5.3), Inches(1.7),
             "Same user, same data, same AI site —\nenforcement differs by who they are right now.\n\nIRM score escalates from events.\nDLP responds in real time.\nNo admin intervention.",
             size=12, bold=True, color=NAVY)


def slide_act5_sentinel():
    s = add_slide()
    add_header(s, "Act 5  |  Sentinel", "Unified SIEM pane + correlated signals")

    # connectors row
    add_text(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(0.5),
             "Three connectors", size=16, bold=True, color=NAVY)
    conns = [
        ("Defender XDR", "DLP alerts via XDR pipeline", BLUE),
        ("Office IRM", "High-sev IRM alerts", TEAL),
        ("Office 365", "Unified audit, Copilot activity", SLATE),
    ]
    for i, (t, body, color) in enumerate(conns):
        left = Inches(0.6 + i * 2.05)
        add_rounded(s, left, Inches(2.5), Inches(1.95), Inches(1.2), WHITE, line=LIGHT)
        add_rect(s, left, Inches(2.5), Inches(0.15), Inches(1.2), color)
        add_text(s, left + Inches(0.25), Inches(2.6), Inches(1.6), Inches(0.4),
                 t, size=12, bold=True, color=NAVY)
        add_text(s, left + Inches(0.25), Inches(2.95), Inches(1.6), Inches(0.7),
                 body, size=10, color=SLATE)

    # rules list
    add_text(s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(0.5),
             "7 analytics rules", size=16, bold=True, color=NAVY)
    rules = [
        "HighSevDLP — high-severity DLP alerts",
        "IRMHighSev — Insider Risk escalations",
        "LabelDowngrade — pre-exfiltration label stripping",
        "MassDownloadAfterDLP — cross-table (DLP + mass download, 4h)",
        "CopilotDLPPromptBlock  (AI-specific)",
        "ShadowAIPasteUpload  (AI-specific)",
        "RiskyAIUsageCorrel  (the key rule — next slide)",
    ]
    add_bullets(s, Inches(7.0), Inches(2.5), Inches(5.7), Inches(4.0), rules, size=12)

    # bottom row: workbooks + playbook
    add_rounded(s, Inches(0.6), Inches(4.0), Inches(6.2), Inches(2.7), WHITE, line=LIGHT)
    add_text(s, Inches(0.8), Inches(4.15), Inches(6.0), Inches(0.5),
             "2 workbooks + 1 playbook", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.8), Inches(4.7), Inches(6.0), Inches(2.0), [
        "Purview Signals — DLP volume, IRM sev, label activity",
        "AI Risk Signals — Copilot blocks, Shadow AI by site, cross-signal users",
        "IRM auto-triage Logic App — MSI, no credentials, sub-second comment",
    ], size=11)


def slide_correlation():
    s = add_slide()
    add_header(s, "Cross-signal correlation", "PVAISec-RiskyAIUsageCorrel — the key rule", kicker_color=AMBER)

    # two input panels
    add_rounded(s, Inches(0.6), Inches(2.0), Inches(4.5), Inches(2.0), WHITE, line=LIGHT)
    add_rect(s, Inches(0.6), Inches(2.0), Inches(4.5), Inches(0.5), BLUE)
    add_text(s, Inches(0.8), Inches(2.1), Inches(4.2), Inches(0.4),
             "Risky AI IRM alerts", size=14, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.7), Inches(4.2), Inches(1.2),
             "Users flagged for Copilot prompt injection\nor protected material access.\nLast 4 hours.",
             size=12, color=SLATE)

    add_rounded(s, Inches(0.6), Inches(4.3), Inches(4.5), Inches(2.0), WHITE, line=LIGHT)
    add_rect(s, Inches(0.6), Inches(4.3), Inches(4.5), Inches(0.5), TEAL)
    add_text(s, Inches(0.8), Inches(4.4), Inches(4.2), Inches(0.4),
             "DLP blocks on AI surfaces", size=14, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(5.0), Inches(4.2), Inches(1.2),
             "Copilot DLP blocks + Shadow AI\npaste/upload blocks.\nSame 4-hour window.",
             size=12, color=SLATE)

    # arrow
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.3), Inches(3.5), Inches(1.5), Inches(1.0))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = AMBER
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    add_text(s, Inches(5.3), Inches(3.6), Inches(1.5), Inches(0.8),
             "JOIN", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # output panel
    add_rounded(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.3), NAVY)
    add_text(s, Inches(7.2), Inches(2.15), Inches(5.3), Inches(0.5),
             "Elevated incident", size=16, bold=True, color=TEAL)
    add_text(s, Inches(7.2), Inches(2.7), Inches(5.3), Inches(3.5),
             "Two soft signals → one hard incident.\n\nAuto-triage playbook enriches\nthe incident with triage guidance.\n\nSOC sees a single item:\n'User combining IRM + DLP on AI surfaces.'\n\nThe investigation-worthy one.",
             size=13, color=WHITE)

    add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.4),
             "This is the money rule. Everything else is feeder data.",
             size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def slide_asymmetry():
    s = add_slide()
    add_header(s, "The asymmetry", "Same detection, different outcome — by design")

    # Two columns
    add_rounded(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(4.5), WHITE, line=LIGHT)
    add_rect(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(0.7), BLUE)
    add_text(s, Inches(0.8), Inches(2.12), Inches(5.6), Inches(0.5),
             "Copilot  (sanctioned)", size=16, bold=True, color=WHITE)
    add_bullets(s, Inches(0.8), Inches(2.85), Inches(5.5), Inches(3.5), [
        "Runs inside your compliance boundary",
        "DLP guides users to the right path",
        "Block with a reason the user understands",
        "Labels travel with content into responses",
        "Lower friction — it's meant to be used",
    ], size=13)

    add_rounded(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(4.5), WHITE, line=LIGHT)
    add_rect(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(0.7), RED)
    add_text(s, Inches(7.0), Inches(2.12), Inches(5.6), Inches(0.5),
             "External AI  (unsanctioned)", size=16, bold=True, color=WHITE)
    add_bullets(s, Inches(7.0), Inches(2.85), Inches(5.5), Inches(3.5), [
        "Outside your compliance boundary",
        "DLP hard-blocks the action",
        "Three layers: device, browser, network",
        "Tightens automatically with IRM score",
        "Higher friction — that's the point",
    ], size=13)

    add_text(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.4),
             "The friction asymmetry is the behavior change lever.",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_prereqs():
    s = add_slide()
    add_header(s, "Prerequisites", "What you need before deploying")
    left_items = [
        "Microsoft 365 E5  (or E5 Compliance)",
        "Microsoft 365 Copilot licenses for demo users",
        "Azure subscription  (Owner or Contributor + UAA)",
        "Defender for Endpoint on ≥1 test device",
        "Azure CLI signed in  (az login)",
        "PowerShell 7+",
    ]
    right_items = [
        "ExchangeOnlineManagement 3.0+",
        "Microsoft.Graph SDK  (Users, Groups, Auth)",
        "Roles: Compliance Admin, User Admin, eDiscovery Admin",
        "Providers: Microsoft.OperationalInsights, Microsoft.SecurityInsights",
        "TenantId + SubscriptionId on hand",
        "4 hour window for DLP propagation before demo",
    ]
    add_text(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(0.5),
             "Licensing + access", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(4.2), left_items, size=14)

    add_text(s, Inches(7.0), Inches(1.95), Inches(6.0), Inches(0.5),
             "Tooling + permissions", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(2.5), Inches(6.0), Inches(4.2), right_items, size=14)


def slide_deploy_cmds():
    s = add_slide()
    add_header(s, "Deploy", "One command, end-to-end")

    # command block
    add_rounded(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.2), NAVY)
    add_text(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.4),
             "DEPLOY", size=12, bold=True, color=TEAL, font="Consolas")
    add_text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.6),
             "az login\naz account set --subscription <sub>\n./Deploy-Lab.ps1 -Cloud commercial -LabProfile ai-security `\n    -TenantId <tenant-guid> -SubscriptionId <sub-guid>",
             size=14, color=WHITE, font="Consolas")

    # readiness block
    add_rounded(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.4), LIGHT)
    add_text(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.4),
             "POST-DEPLOY READINESS", size=12, bold=True, color=BLUE, font="Consolas")
    add_text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(1.8),
             "./scripts/Test-CopilotDlpReady.ps1  -ConfigPath ./configs/commercial/ai-security-demo.json\n"
             "./scripts/Test-ShadowAiReady.ps1    -ConfigPath ./configs/commercial/ai-security-demo.json\n"
             "./scripts/Test-SentinelReady.ps1    -ConfigPath ./configs/commercial/ai-security-demo.json\n"
             "./scripts/Set-ShadowAiEndpointDlpDomains.ps1 -ConfigPath ./configs/... -Apply",
             size=12, color=SLATE, font="Consolas")

    add_text(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4),
             "Deploy ~20-25 min  |  DLP propagation ~4h  |  Sentinel connector data flow ~60 min",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)


def slide_demo_flow():
    s = add_slide()
    add_header(s, "Demo flow", "45-75 minutes hands-on  |  30 minutes narrative")
    acts = [
        ("Open", "2 min", "Why integrated — three stories, one system", BLUE),
        ("Discovery", "3 min", "Audit searches + AI Risk Signals workbook", TEAL),
        ("Copilot DLP", "5-10 min", "Prompt SIT Block + Labeled Content Block", BLUE),
        ("Shadow AI", "10-15 min", "Device + browser + network live demo", TEAL),
        ("Insider Risk", "5 min", "Adaptive tier escalation walkthrough", AMBER),
        ("Sentinel", "10-15 min", "Rules, workbooks, cross-signal correlation", NAVY),
        ("Investigate + teardown", "5 min", "Auto-triage playbook + how to tear down", SLATE),
        ("DSPM for AI", "5 min", "Optional posture follow-up", MUTED),
    ]
    y = Inches(1.95)
    for i, (t, dur, body, color) in enumerate(acts):
        row = i // 2
        col = i % 2
        left = Inches(0.6 + col * 6.1)
        top = y + row * Inches(1.25)
        add_rounded(s, left, top, Inches(5.95), Inches(1.1), WHITE, line=LIGHT)
        add_rect(s, left, top, Inches(0.15), Inches(1.1), color)
        add_text(s, left + Inches(0.3), top + Inches(0.1), Inches(4.0), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, left + Inches(4.3), top + Inches(0.1), Inches(1.5), Inches(0.4),
                 dur, size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
        add_text(s, left + Inches(0.3), top + Inches(0.55), Inches(5.5), Inches(0.5),
                 body, size=11, color=SLATE)


def slide_references():
    s = add_slide()
    add_header(s, "References", "Microsoft Learn + lab artifacts", kicker_color=BLUE)

    add_text(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(0.5),
             "Microsoft Learn", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(4.0), [
        "DLP for Microsoft 365 Copilot",
        "Shadow AI deployment guide",
        "Sentinel + Purview integration",
        "DSPM for AI",
        "Microsoft Sentinel in the Defender portal",
    ], size=13)

    add_text(s, Inches(7.0), Inches(1.95), Inches(6.0), Inches(0.5),
             "In this repo", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(2.5), Inches(6.0), Inches(4.0), [
        "profiles/commercial/ai-security/README.md",
        "profiles/commercial/ai-security/talk-track.md",
        "profiles/commercial/ai-security/RUNBOOK.md",
        "profiles/commercial/ai-security/demo-scenarios.json",
        "configs/commercial/ai-security-demo.json",
    ], size=12)


def slide_close():
    s = add_slide()
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    add_rect(s, Inches(0.6), Inches(2.2), Inches(0.2), Inches(3.0), TEAL)
    add_text(s, Inches(1.0), Inches(2.0), Inches(12.0), Inches(0.5),
             "THE RECAP", size=14, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(2.4), Inches(12.0), Inches(1.0),
             "Not four tools. One system.", size=34, bold=True, color=WHITE)
    add_bullets(s, Inches(1.0), Inches(3.6), Inches(12.0), Inches(3.0), [
        "Sanctioned AI — Copilot DLP at prompt + label level",
        "Unsanctioned AI — device, browser, network",
        "Insider Risk — the adaptive bridge",
        "Sentinel — unified SIEM, 7 rules, auto-triage",
        "The loop — signals reinforce each other, no admin intervention",
    ], size=16, color=WHITE, bold=False)
    add_text(s, Inches(1.0), Inches(6.7), Inches(12.0), Inches(0.5),
             "Deployed programmatically in 20-25 min from one config file.",
             size=14, color=TEAL, bold=True)


# ---------- Build ----------
slide_title()
slide_problem()
slide_thesis()
slide_loop()
slide_deploy_overview()
slide_act1()
slide_act2_copilot()
slide_act3_shadow()
slide_asymmetry()
slide_act4_irm()
slide_act5_sentinel()
slide_correlation()
slide_demo_flow()
slide_prereqs()
slide_deploy_cmds()
slide_references()
slide_close()

# apply footers with correct totals
total = len(prs.slides)
for i, sl in enumerate(prs.slides, start=1):
    if i == 1 or i == total:
        continue
    add_footer(sl, i, total)

out = "/Users/charlesshea/repos/purview-lab-deployer/profiles/commercial/ai-security/ai-security-lab.pptx"
prs.save(out)
print(f"wrote {out}  ({total} slides)")
